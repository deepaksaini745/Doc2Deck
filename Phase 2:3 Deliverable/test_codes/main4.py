import os
import json
import urllib.request
import urllib.error
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from llama_cloud_services import LlamaParse
from dotenv import load_dotenv
from prompt_templates import EXTRACT_TOPICS_MARKERS_TEMPLATE, GENERATE_SLIDE_CONTENT_TEMPLATE
from fuzzywuzzy import fuzz
import random

# Load environment variables from .env file
load_dotenv()

LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

base_dir = os.path.dirname(os.path.abspath(__file__))
intermediate_dir = os.path.join(base_dir, "intermediate")
os.makedirs(intermediate_dir, exist_ok=True)

def make_api_call(api_key, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': api_key,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except Exception as e:
        print(f"API error: {e}")
    return None

def extract_document_data(doc_path):
    parser = LlamaParse(
        api_key=LLAMA_CLOUD_API_KEY,
        language="en",
        verbose=True
    )

    result = parser.parse(doc_path)

    text = "\n".join(doc.text for doc in result.get_text_documents(split_by_page=False))

    images = result.get_image_documents(
        include_screenshot_images=True,
        include_object_images=True,
        image_download_dir="./images"
    )
    image_data = []
    for img in images:
        if hasattr(img, "file_path"):
            context_text = "No description available"
            for page in result.pages:
                if hasattr(page, "images") and page.images:
                    for page_image in page.images:
                        if hasattr(page_image, "file_path") and page_image.file_path == img.file_path:
                            context_parts = []
                            if hasattr(page, "layout") and page.layout:
                                for block in page.layout:
                                    if hasattr(block, "type") and block.type in ["heading", "paragraph"]:
                                        if hasattr(block, "text") and block.text.strip():
                                            context_parts.append(block.text.strip())
                            if context_parts:
                                context_text = " ".join(context_parts)
                            break
            image_data.append({"path": img.file_path, "context": context_text})

    print("\n========= Extracted Image Contexts =========")
    if not image_data:
        print("No images extracted from the document.")
    else:
        for image in image_data:
            print(f"Image path: {image['path']}")
            print(f"Context: {image['context']}")
            print("-" * 50)
    print("==========================================\n")

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    return {"text": text, "images": image_data, "tables": tables}

def find_marker_position(text, marker):
    position = text.find(marker)
    if position != -1:
        return position

    words = text.split()
    marker_words = marker.split()
    best_match = -1
    best_score = 0

    for i in range(len(words) - len(marker_words) + 1):
        window = ' '.join(words[i:i + len(marker_words)])
        score = fuzz.ratio(window.lower(), marker.lower())
        if score > best_score:
            best_score = score
            best_match = i

    if best_score > 70:
        return len(' '.join(words[:best_match]))
    else:
        print(f"Warning: Could not find good position for marker: '{marker[:30]}...' (Best score: {best_score})")
        return -1

def generate_topics_and_segments(text):
    extract_prompt = EXTRACT_TOPICS_MARKERS_TEMPLATE.replace("{{content}}", text)
    topics_and_markers = make_api_call(ANTHROPIC_API_KEY, extract_prompt)
    if not topics_and_markers:
        raise ValueError("Failed to extract topics and markers.")

    topics = []
    lines = topics_and_markers.strip().split('\n')
    current_topic = None
    current_marker = ""
    for line in lines:
        if line.startswith('**') and line.endswith('*'):
            if current_topic:
                topics.append({"topic": current_topic, "marker": current_marker.strip()})
            current_topic = line.strip('*').strip()
            current_marker = ""
        else:
            current_marker += line + " "
    if current_topic:
        topics.append({"topic": current_topic, "marker": current_marker.strip()})

    return topics

def segment_content(text, topics):
    segments = []
    for i, item in enumerate(topics):
        start = find_marker_position(text, item['marker'])
        if start == -1:
            continue
        if i + 1 < len(topics):
            end = find_marker_position(text, topics[i + 1]['marker'])
            end = end if end != -1 else len(text)
        else:
            end = len(text)

        if start < end:
            segment = text[start:end].strip()
            segments.append({"topic": item['topic'], "content": segment})

    return segments

def generate_slide_contents(segments):
    slides = []
    for segment in segments:
        prompt = GENERATE_SLIDE_CONTENT_TEMPLATE.replace("{{topic}}", segment['topic']).replace("{{contentSegment}}", segment['content'])
        slide_text = make_api_call(ANTHROPIC_API_KEY, prompt)
        if slide_text:
            slides.append({"topic": segment['topic'], "text": slide_text})
        else:
            print(f"Warning: No slide content generated for topic: {segment['topic']}")
    return slides

def create_ppt(slides, image_data, output_ppt_path, template_path):
    ppt = Presentation(template_path)

    if len(ppt.slides) > 0:
        rId = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(rId)
        del ppt.slides._sldIdLst[0]

    for slide_data in slides:
        lines = slide_data["text"].strip().split("\n")

        current_slide = None

        for line in lines:
            line = line.strip().replace("**", "")
            if not line:
                continue

            if not line.startswith("•") and not line.startswith("-"):
                current_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                current_slide.shapes.title.text = line

                text_box = current_slide.placeholders[1]
                text_frame = text_box.text_frame
                text_frame.clear()

                if image_data:
                    image_choice = random.choice(image_data)
                    image_path = image_choice["path"]
                    try:
                        current_slide.shapes.add_picture(image_path, Inches(5), Inches(1.5), width=Inches(3))
                    except Exception as e:
                        print(f"Error adding image {image_path}: {e}")

            elif current_slide:
                p = current_slide.placeholders[1].text_frame.add_paragraph()
                p.text = line.replace("- ", "• ").replace("• ", "• ").strip()
                p.font.size = Pt(18)

    # Add Thank You slide at the end
    thank_you_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    thank_you_slide.shapes.title.text = "Thank You!"

    thank_you_textbox = thank_you_slide.placeholders[1]
    thank_you_text_frame = thank_you_textbox.text_frame
    thank_you_text_frame.clear()

    p = thank_you_text_frame.add_paragraph()
    p.text = "We appreciate your attention. Looking forward to your questions."
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.LEFT

    # Optional: Add an image if available
    if image_data:
        image_choice = random.choice(image_data)
        image_path = image_choice["path"]
        try:
            thank_you_slide.shapes.add_picture(image_path, Inches(4), Inches(2), width=Inches(3))
        except Exception as e:
            print(f"Error adding image to Thank You slide: {e}")

    ppt.save(output_ppt_path)
    print(f"Presentation saved as {output_ppt_path}")

def main():
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    output_ppt_path = os.path.join(base_dir, "output", "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")

    document_data = extract_document_data(doc_path)
    topics = generate_topics_and_segments(document_data["text"])
    segments = segment_content(document_data["text"], topics)
    slides = generate_slide_contents(segments)

    create_ppt(slides, document_data["images"], output_ppt_path, template_path)

if __name__ == "__main__":
    main()