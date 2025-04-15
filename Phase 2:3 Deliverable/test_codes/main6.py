import os
import json
import urllib.request
import urllib.error
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from llama_cloud_services import LlamaParse
from dotenv import load_dotenv
from prompt_templates import ENRICH_PRESENTATION_PROMPT, EXTRACT_TOPICS_MARKERS_TEMPLATE, GENERATE_SLIDE_CONTENT_TEMPLATE
from fuzzywuzzy import fuzz
import random

# Load environment variables from .env file
load_dotenv()

LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

base_dir = os.path.dirname(os.path.abspath(__file__))
intermediate_dir = os.path.join(base_dir, "intermediate")
os.makedirs(intermediate_dir, exist_ok=True)

def safe_prompt(template, variables):
    # First, do the replacements properly
    for key, value in variables.items():
        template = template.replace(f"{{{{{key}}}}}", value)
    # Now escape any leftover braces in the replaced content
    template = template.replace("{", "{{").replace("}", "}}")
    return template


def make_api_call(api_key, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': api_key,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    print("\n====== API CALL PROMPT ======")
    print(content)
    print("====== END PROMPT ======\n")

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            parsed_response = json.loads(output)['content'][0]['text']
            print("\n====== API RAW RESPONSE ======")
            print(parsed_response)
            print("====== END RESPONSE ======\n")
            return parsed_response
    except Exception as e:
        print(f"API error: {e}")
    return None

def extract_document_data(doc_path):
    parser = LlamaParse(
        api_key=LLAMA_CLOUD_API_KEY,
        language="en",
        verbose=True
    )

    result = parser.parse(doc_path)

    text = "\n".join(doc.text for doc in result.get_text_documents(split_by_page=False))

    image_data = []
    for page in result.pages:
        if hasattr(page, "layout"):
            for block in page.layout:
                if getattr(block, "type", None) == "image" and hasattr(block, "image"):
                    if hasattr(block.image, "file_path"):
                        image_data.append({
                            "path": block.image.file_path,
                            "context": "Context not available inline"
                        })

    print("\n========= Extracted Image Contexts =========")
    if not image_data:
        print("No images extracted from the document.")
    else:
        for image in image_data:
            print(f"Image path: {image['path']}")
            print(f"Context: {image['context']}")
            print("-" * 50)
    print("==========================================\n")

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    print("\n====== DEBUG: Extracted Tables ======")
    if not tables:
        print("No tables extracted from the document.")
    else:
        for table in tables:
            for row in table:
                print(row)
            print("-" * 50)
    print("====================================\n")

    print("\n====== DEBUG: Extracted Text ======")
    print(text[:2000])
    print("... (text truncated)")
    print("====================================\n")

    return {"text": text, "images": image_data, "tables": tables}

def find_marker_position(text, marker):
    position = text.find(marker)
    if position != -1:
        return position

    words = text.split()
    marker_words = marker.split()
    best_match = -1
    best_score = 0

    for i in range(len(words) - len(marker_words) + 1):
        window = ' '.join(words[i:i + len(marker_words)])
        score = fuzz.ratio(window.lower(), marker.lower())
        if score > best_score:
            best_score = score
            best_match = i

    if best_score > 70:
        return len(' '.join(words[:best_match]))
    else:
        print(f"Warning: Could not find good position for marker: '{marker[:30]}...' (Best score: {best_score})")
        return -1

def generate_topics_and_segments(text):
    extract_prompt = safe_prompt(EXTRACT_TOPICS_MARKERS_TEMPLATE, {"content": text})
    topics_and_markers = make_api_call(ANTHROPIC_API_KEY, extract_prompt)
    if not topics_and_markers:
        raise ValueError("Failed to extract topics and markers.")

    topics = []
    lines = topics_and_markers.strip().split('\n')
    current_topic = None
    current_marker = ""
    for line in lines:
        if line.startswith('**') and line.endswith('**'):
            if current_topic:
                topics.append({"topic": current_topic, "marker": current_marker.strip()})
            current_topic = line.strip('*').strip()
            current_marker = ""
        else:
            current_marker += line + " "
    if current_topic:
        topics.append({"topic": current_topic, "marker": current_marker.strip()})

    return topics

def segment_content(text, topics):
    segments = []
    for i, item in enumerate(topics):
        start = find_marker_position(text, item['marker'])
        if start == -1:
            continue
        if i + 1 < len(topics):
            end = find_marker_position(text, topics[i + 1]['marker'])
            end = end if end != -1 else len(text)
        else:
            end = len(text)

        if start < end:
            segment = text[start:end].strip()
            segments.append({"topic": item['topic'], "content": segment})

    return segments

def generate_slide_contents(segments):
    slides = []
    for segment in segments:
        prompt = safe_prompt(GENERATE_SLIDE_CONTENT_TEMPLATE, {"topic": segment['topic'], "contentSegment": segment['content']})
        slide_text = make_api_call(ANTHROPIC_API_KEY, prompt)
        if slide_text:
            slides.append({"topic": segment['topic'], "text": slide_text})
        else:
            print(f"Warning: No slide content generated for topic: {segment['topic']}")
    return slides

def add_table_slide(ppt, table_data):
    if not table_data:
        return

    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "Summary Table"

    rows, cols = len(table_data), len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(0.8)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    for i, row in enumerate(table_data):
        for j, cell in enumerate(row):
            table_shape.cell(i, j).text = cell
            table_shape.cell(i, j).text_frame.paragraphs[0].font.size = Pt(12)

    print("Added summary table slide.")

def create_ppt(slides, image_data, table_data, output_ppt_path, template_path):
    ppt = Presentation(template_path)

    if len(ppt.slides) > 0:
        rId = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(rId)
        del ppt.slides._sldIdLst[0]

    for slide_data in slides:
        lines = slide_data["text"].strip().split("\n")
        current_slide = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            if not line.startswith("•") and not line.startswith("-"):
                current_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                current_slide.shapes.title.text = line.replace("**", "")

                text_box = current_slide.placeholders[1]
                text_frame = text_box.text_frame
                text_frame.clear()

                if image_data:
                    image_choice = random.choice(image_data)
                    image_path = image_choice["path"]
                    try:
                        current_slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(5.5))
                    except Exception as e:
                        print(f"Error adding image {image_path}: {e}")

            elif current_slide:
                p = current_slide.placeholders[1].text_frame.add_paragraph()
                p.text = line.replace("- ", "• ").replace("• ", "• ").strip()
                p.font.size = Pt(18)

    add_table_slide(ppt, table_data)

    thank_you_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
    thank_you_slide.shapes.title.text = "Thank You!"

    text_box = thank_you_slide.placeholders[1]
    text_frame = text_box.text_frame
    text_frame.clear()

    p1 = text_frame.add_paragraph()
    p1.text = "We appreciate your attention."
    p1.font.size = Pt(24)

    p2 = text_frame.add_paragraph()
    p2.text = "Looking forward to your questions and discussions!"
    p2.font.size = Pt(18)

    if image_data:
        image_choice = random.choice(image_data)
        image_path = image_choice["path"]
        try:
            thank_you_slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(10), height=Inches(5.5))
        except Exception as e:
            print(f"Error adding background image to Thank You slide: {e}")

    ppt.save(output_ppt_path)
    print(f"Presentation saved as {output_ppt_path}")

def main():
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    output_ppt_path = os.path.join(base_dir, "output", "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")

    document_data = extract_document_data(doc_path)
    topics = generate_topics_and_segments(document_data["text"])
    segments = segment_content(document_data["text"], topics)
    slides = generate_slide_contents(segments)

    create_ppt(slides, document_data["images"], document_data["tables"], output_ppt_path, template_path)

if __name__ == "__main__":
    main()
