import os
import json
import urllib.request
import urllib.error
from pptx import Presentation
from pptx.util import Inches
from llama_cloud_services import LlamaParse
from dotenv import load_dotenv
from prompt_templates import ENRICH_PRESENTATION_PROMPT, EXTRACT_TOPICS_MARKERS_TEMPLATE, GENERATE_SLIDE_CONTENT_TEMPLATE  # Import all your templates

# Load environment variables from .env file
load_dotenv()

# Retrieve API keys from the environment
LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

def extract_document_data(doc_path):
    parser = LlamaParse(
        api_key=LLAMA_CLOUD_API_KEY,
        language="en",
        verbose=True
    )

    result = parser.parse(doc_path)

    text_documents = result.get_text_documents(split_by_page=False)
    text = "\n".join(doc.text for doc in text_documents)

    images = result.get_image_documents(
        include_screenshot_images=True,
        include_object_images=True,
        image_download_dir="./images"
    )
    image_paths = [img.file_path for img in images if hasattr(img, "file_path")]

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    return {"text": text, "images": image_paths, "tables": tables}

def make_api_call(api_key, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': api_key,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    print(f"Attempting to access: {url}")
    print(f"With headers: {headers}")

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except urllib.error.HTTPError as e:
        print(f"HTTP error: {e.code} - {e.reason}")
        print(f"Response body: {e.read().decode()}")
    except urllib.error.URLError as e:
        print(f"URL error: {e}")
    except Exception as e:
        print(f"General error: {e}")
    return None

def enrich_with_claude(document_data):
    # Using ENRICH_PRESENTATION_PROMPT from prompt_templates
    formatted_prompt = ENRICH_PRESENTATION_PROMPT.format(
        text_content=document_data['text'],
        image_paths=document_data['images'],
        table_data=json.dumps(document_data['tables'], indent=2)
    )

    response = make_api_call(ANTHROPIC_API_KEY, formatted_prompt)

    if response is None:
        raise ValueError("Claude API call failed. No response received.")

    try:
        structured_response = json.loads(response)
    except json.JSONDecodeError:
        raise ValueError("Claude response is not valid JSON. Response was:\n" + response)

    return structured_response

def create_ppt_from_claude(claude_slides, ppt_path="generated_deck.pptx", template_path=None):
    if template_path and os.path.exists(template_path):
        ppt = Presentation(template_path)
    else:
        ppt = Presentation()

    for slide_data in claude_slides:
        # Use layout 1 if available (Title and Content)
        # Clean up all placeholders at the beginning to avoid template residuals
        if len(ppt.slide_layouts) > 1:
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[0])

        # Clean up placeholders early
        for shape in list(slide.placeholders):
            try:
                if not shape.has_text_frame or not shape.text_frame.text.strip():
                    sp = shape
                    sp.element.getparent().remove(sp.element)
            except Exception:
                pass
        # Slide title
        if slide_data.get("title"):
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
            else:
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
                title_shape.text = slide_data["title"]

        # Slide text
        if slide_data.get("text"):
            # Try to use placeholder for body text
            text_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # BODY placeholder
                    text_placeholder = shape
                    break
            if text_placeholder:
                text_frame = text_placeholder.text_frame
                text_frame.clear()  # Clear any default text from placeholder

                lines = slide_data["text"].split("\n")
                clean_lines = [line.lstrip("-• ").strip() for line in lines if line.strip() != ""]

                if clean_lines:
                    text_frame.text = clean_lines[0]
                    for line in clean_lines[1:]:
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
                        p.space_before = Inches(0.05)
                        p.space_after = Inches(0.05)
                        p.text_frame.auto_size = True
                        p.text_frame.word_wrap = True

                    # Auto-adjust font size if too many bullets
                    if len(clean_lines) > 5:
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Inches(0.2)

                    # Improve bullet spacing
                    for paragraph in text_frame.paragraphs:
                        paragraph.space_before = Inches(0.05)
                        paragraph.space_after = Inches(0.05)

                # Remove unused placeholders (clean up empty icons)
                for shape in slide.placeholders:
                    if not shape.has_text_frame or not shape.text_frame.text.strip():
                        try:
                            sp = shape
                            sp.element.getparent().remove(sp.element)
                        except Exception:
                            pass
            else:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(3))
                text_frame = textbox.text_frame
                lines = slide_data["text"].split("\n")
                clean_lines = [line.lstrip("-• ").strip() for line in lines if line.strip() != ""]

                if clean_lines:
                    text_frame.auto_size = True
                    text_frame.word_wrap = True

                    for line in clean_lines:
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
                        p.space_before = Inches(0.05)
                        p.space_after = Inches(0.05)

        # Slide image
        if slide_data.get("image") and os.path.exists(slide_data["image"]):
            # Center image horizontally below text
            img_path = slide_data["image"]
            pic = slide.shapes.add_picture(img_path, Inches(1), Inches(4.5), width=Inches(6))

        # Slide table
        if slide_data.get("table"):
            table_data = slide_data["table"]
            if table_data and isinstance(table_data, list) and table_data[0]:
                rows, cols = len(table_data), len(table_data[0])
                tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.5), Inches(8), Inches(2)).table
                for i, row in enumerate(table_data):
                    for j, cell in enumerate(row):
                        cell_obj = tbl_shape.cell(i, j)
                        cell_obj.text = str(cell)
                        # Align text centrally
                        cell_obj.text_frame.paragraphs[0].alignment = 1  # Center

                # Adjust column widths evenly
                total_width = Inches(8)
                for col in range(cols):
                    tbl_shape.columns[col].width = int(total_width / cols)

    ppt.save(ppt_path)
    print(f"Presentation saved as {ppt_path}")

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    document_data = extract_document_data(doc_path)

    intermediate_dir = os.path.join(base_dir, "intermediate")
    os.makedirs(intermediate_dir, exist_ok=True)

    extracted_data_path = os.path.join(intermediate_dir, "extracted_data.json")
    with open(extracted_data_path, "w", encoding="utf-8") as json_file:
        json.dump(document_data, json_file, ensure_ascii=False, indent=4)

    claude_response = enrich_with_claude(document_data)

    claude_response_path = os.path.join(intermediate_dir, "claude_structured_response.json")
    with open(claude_response_path, "w", encoding="utf-8") as json_file:
        json.dump(claude_response, json_file, ensure_ascii=False, indent=4)

    output_dir = os.path.join(base_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    output_ppt_path = os.path.join(output_dir, "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")
    create_ppt_from_claude(claude_response, output_ppt_path, template_path)

if __name__ == "__main__":
    main()
