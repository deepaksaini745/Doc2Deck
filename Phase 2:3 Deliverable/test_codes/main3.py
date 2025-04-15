import os
import json
import urllib.request
import urllib.error
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from llama_cloud_services import LlamaParse
from dotenv import load_dotenv
from prompt_templates import ENRICH_PRESENTATION_PROMPT, EXTRACT_TOPICS_MARKERS_TEMPLATE, GENERATE_SLIDE_CONTENT_TEMPLATE  # Import all your templates
from fuzzywuzzy import fuzz

# Load environment variables from .env file
load_dotenv()

# Retrieve API keys from the environment
LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

def extract_document_data(doc_path):
    parser = LlamaParse(
        api_key=LLAMA_CLOUD_API_KEY,
        language="en",
        verbose=True
    )

    result = parser.parse(doc_path)

    text_documents = result.get_text_documents(split_by_page=False)
    text = "\n".join(doc.text for doc in text_documents)

    images = result.get_image_documents(
        include_screenshot_images=True,
        include_object_images=True,
        image_download_dir="./images"
    )
    image_data = []
    for img in images:
        if hasattr(img, "file_path"):
            context_text = "No description available"
            for page in result.pages:
                if hasattr(page, "images") and page.images:
                    for page_image in page.images:
                        if hasattr(page_image, "file_path") and page_image.file_path == img.file_path:
                            page_context_parts = []
                            if hasattr(page, "layout") and page.layout:
                                for idx, block in enumerate(page.layout):
                                    if hasattr(block, "type") and block.type in ["heading", "paragraph"]:
                                        if hasattr(block, "text") and block.text.strip():
                                            page_context_parts.append(block.text.strip())
                            if page_context_parts:
                                context_text = " ".join(page_context_parts)
                            break
            image_data.append({
                "path": img.file_path,
                "context": context_text
            })

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    return {"text": text, "images": image_data, "tables": tables}

def make_api_call(api_key, content):
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': api_key,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    print(f"Attempting to access: {url}")
    print(f"With headers: {headers}")

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except urllib.error.HTTPError as e:
        print(f"HTTP error: {e.code} - {e.reason}")
        print(f"Response body: {e.read().decode()}")
    except urllib.error.URLError as e:
        print(f"URL error: {e}")
    except Exception as e:
        print(f"General error: {e}")
    return None

def enrich_with_claude(document_data):
    formatted_prompt = ENRICH_PRESENTATION_PROMPT.format(
        text_content=document_data['text'],
        image_paths=json.dumps(document_data['images'], indent=2),
        table_data=json.dumps(document_data['tables'], indent=2)
    )

    response = make_api_call(ANTHROPIC_API_KEY, formatted_prompt)

    if response is None:
        raise ValueError("Claude API call failed. No response received.")

    try:
        structured_response = json.loads(response)
    except json.JSONDecodeError:
        raise ValueError("Claude response is not valid JSON. Response was:\n" + response)

    return structured_response

def find_best_image_for_slide(slide, images):
    slide_text = f"{slide.get('title', '')} {slide.get('text', '')}"
    best_score = 0
    best_image = None
    for image in images:
        score = fuzz.partial_ratio(slide_text.lower(), image['context'].lower())
        if score > best_score:
            best_score = score
            best_image = image['path']
    return best_image

def create_ppt_from_claude(claude_slides, ppt_path="generated_deck.pptx", template_path=None):
    images_dir = "./images"
    available_images = []
    if os.path.exists(images_dir):
        for img_file in os.listdir(images_dir):
            if img_file.lower().endswith((".png", ".jpg", ".jpeg")):
                available_images.append(os.path.join(images_dir, img_file))

    for slide_data in claude_slides:
        if (not slide_data.get("image") or not os.path.exists(slide_data.get("image"))) and available_images:
            best_image = find_best_image_for_slide(slide_data, [{"path": img, "context": ""} for img in available_images])
            slide_data["image"] = best_image if best_image else available_images[0]

    if template_path and os.path.exists(template_path):
        ppt = Presentation(template_path)
        while ppt.slides:
            rId = ppt.slides._sldIdLst[0].rId
            ppt.part.drop_rel(rId)
            del ppt.slides._sldIdLst[0]
    else:
        ppt = Presentation()

    for slide_data in claude_slides:
        if len(ppt.slide_layouts) > 1:
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        else:
            slide = ppt.slides.add_slide(ppt.slide_layouts[0])

        for shape in list(slide.placeholders):
            try:
                if not shape.has_text_frame or not shape.text_frame.text.strip():
                    sp = shape
                    sp.element.getparent().remove(sp.element)
            except Exception:
                pass

        if slide_data.get("title"):
            if slide.shapes.title:
                title_shape = slide.shapes.title
            else:
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))

            title_shape.text = slide_data["title"]
            if title_shape.text_frame.paragraphs:
                paragraph = title_shape.text_frame.paragraphs[0]
                paragraph.alignment = PP_ALIGN.LEFT
                if paragraph.runs:
                    run = paragraph.runs[0]
                    run.font.bold = True
                    run.font.size = Pt(32)

        if slide_data.get("text"):
            text_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:
                    text_placeholder = shape
                    break

            lines = slide_data["text"].split("\n")
            clean_lines = [line.lstrip("-• ").strip() for line in lines if line.strip() != ""]

            if text_placeholder and clean_lines:
                text_frame = text_placeholder.text_frame
                text_frame.clear()

                for line in clean_lines:
                    p = text_frame.add_paragraph()
                    p.text = line
                    p.level = 0
                    p.bullet = True
                    p.space_before = Inches(0.05)
                    p.space_after = Inches(0.05)

                    for run in p.runs:
                        run.font.size = Pt(18)
                        run.font.name = 'Calibri'

                text_frame.paragraphs = [p for p in text_frame.paragraphs if p.text.strip() != ""]

            else:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(3))
                text_frame = textbox.text_frame

                if clean_lines:
                    text_frame.auto_size = True
                    text_frame.word_wrap = True

                    for line in clean_lines:
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.level = 0
                        p.bullet = True
                        p.space_before = Inches(0.05)
                        p.space_after = Inches(0.05)

        if slide_data.get("image") and os.path.exists(slide_data["image"]):
            img_path = slide_data["image"]
            slide.shapes.add_picture(img_path, Inches(1), Inches(4.5), width=Inches(6))

        if slide_data.get("table"):
            table_data = slide_data["table"]
            if table_data and isinstance(table_data, list) and table_data[0]:
                rows, cols = len(table_data), len(table_data[0])
                tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(4.5), Inches(8), Inches(2)).table
                for i, row in enumerate(table_data):
                    for j, cell in enumerate(row):
                        cell_obj = tbl_shape.cell(i, j)
                        cell_obj.text = str(cell)
                        cell_obj.text_frame.paragraphs[0].alignment = 1

                total_width = Inches(8)
                for col in range(cols):
                    tbl_shape.columns[col].width = int(total_width / cols)

    ppt.save(ppt_path)
    print(f"Presentation saved as {ppt_path}")

def main():
    base_dir = os.path.dirname(os.path.abspath(__file__))
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    document_data = extract_document_data(doc_path)

    intermediate_dir = os.path.join(base_dir, "intermediate")
    os.makedirs(intermediate_dir, exist_ok=True)

    extracted_data_path = os.path.join(intermediate_dir, "extracted_data.json")
    with open(extracted_data_path, "w", encoding="utf-8") as json_file:
        json.dump(document_data, json_file, ensure_ascii=False, indent=4)

    claude_response = enrich_with_claude(document_data)

    claude_response_path = os.path.join(intermediate_dir, "claude_structured_response.json")
    with open(claude_response_path, "w", encoding="utf-8") as json_file:
        json.dump(claude_response, json_file, ensure_ascii=False, indent=4)

    output_dir = os.path.join(base_dir, "output")
    os.makedirs(output_dir, exist_ok=True)

    output_ppt_path = os.path.join(output_dir, "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")
    create_ppt_from_claude(claude_response, output_ppt_path, template_path)

if __name__ == "__main__":
    main()
