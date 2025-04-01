from figure_extractor import extract_figures_from_docx
from figure_extractor import decide_slide_mapping
# from figure_extractor import extract_all_images
import subprocess
import sys
import json
import re
import os
import urllib.request
import ssl
print(ssl.OPENSSL_VERSION)

def install_and_import(package):
    """Install and import the specified package."""
    try:
        __import__(package)
    except ImportError:
        subprocess.check_call([sys.executable, "-m", "pip", "install", package])

# # Check and install the necessary packages
# install_and_import('python_docx')
# install_and_import('python_pptx')
# install_and_import('fuzzywuzzy')
# install_and_import('python_Levenshtein')
# install_and_import('python-dotenv')

from docx import Document
from pptx import Presentation
from dotenv import load_dotenv
from fuzzywuzzy import fuzz
from prompt_templates import EXTRACT_TOPICS_MARKERS_TEMPLATE, GENERATE_SLIDE_CONTENT_TEMPLATE

class MarkerNotFoundError(Exception):
    """Custom exception for when a marker is not found in the content."""
    pass

def makeApiCall(apiKey, content):
    """Make an API call to the Anthropic Claude model."""
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': apiKey,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": content}]
    }).encode('utf-8')

    print(f"Attempting to access: {url}")
    print(f"With headers: {headers}")

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')

    try:
        with urllib.request.urlopen(req) as response:
            output = response.read()
            return json.loads(output)['content'][0]['text']
    except urllib.error.HTTPError as e:
        print(f"HTTP error: {e.code} - {e.reason}")
        print(f"Response body: {e.read().decode()}")
    except urllib.error.URLError as e:
        print(f"URL error: {e}")
    except Exception as e:
        print(f"General error: {e}")
    return None

def py_extractContentFromWordDoc(wordDocFilePath):
    """Extract content from a Word document."""
    try:
        doc = Document(wordDocFilePath)
        content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        py_saveFile(wordDocTextFilePath, content)
        return content
    except Exception as e:
        print(f"Error reading Word document: {e}")
        return None

def py_generatePrompt(promptTemplate, vars):
    """Generate a prompt by replacing variables in the template."""
    def replace_var(match):
        var_name = match.group(1)
        return str(vars.get(var_name, f"{{{{Undefined variable: {var_name}}}}}"))
    
    prompt = re.sub(r'\{\{(\w+)\}\}', replace_var, promptTemplate)
    # print(prompt)
    return prompt

def py_saveFile(filePath, content):
    """Save content to a file."""
    try:
        with open(filePath, 'w', encoding='utf-8') as file:
            file.write(content)
    except Exception as e:
        print(f"Error saving file: {e}")

def py_convertTextToJson(topicsAndMarkers):
    """Convert text output to JSON format."""
    lines = topicsAndMarkers.strip().split('\n')
    result = []
    current_topic = None
    current_marker = ""
    for line in lines:
        if line.startswith('**') and line.endswith('*'):
            if current_topic:
                result.append({"topic": current_topic, "marker": current_marker.strip()})
            current_topic = line.strip('*').strip()
            current_marker = ""
        else:
            current_marker += line + " "
    if current_topic:  # Add the last topic
        result.append({"topic": current_topic, "marker": current_marker.strip()})
    return json.dumps(result)

def py_getMarkerPosition(marker, content, topic):
    """Get the position of a marker in the content."""
    # First, try for an exact match
    exact_match = content.find(marker)
    if exact_match != -1:
        return exact_match

    # If exact match fails, use fuzzy matching
    best_match_position = -1
    best_match_ratio = 0
    marker_words = marker.split()
    words = content.split()
    
    for i in range(len(words) - len(marker_words) + 1):
        chunk = ' '.join(words[i:i + len(marker_words)])
        ratio = fuzz.ratio(marker.lower(), chunk.lower())
        
        if ratio > best_match_ratio:
            best_match_ratio = ratio
            best_match_position = i  # Store the start index of the match
            
        if ratio > 99:  # Good enough match, break early
            break
    
    if best_match_position == -1:
        content_preview = content[:100] + "..." if len(content) > 100 else content
        raise MarkerNotFoundError(f"Could not find a good match for topic '{topic}' with marker '{marker}'. Content preview: {content_preview}")
    
    # Return the start position of the best matching chunk
    return len(' '.join(words[:best_match_position]))

def py_extractContentSegment(startPosition, endPosition, content):
    """Extract a segment of content between two positions."""
    return content[startPosition:endPosition]

import re
from pptx import Presentation

def py_generatePPT(slides, template_path):
    """Generate a PowerPoint presentation from slide content using a template."""
    prs = Presentation(template_path)
    current_slide = None

    # Regular expression to match "Slide X:" or "Slide X Title:"
    slide_prefix_pattern = re.compile(r'^Slide \d+( Title)?:', re.IGNORECASE)
    # Regular expression to match "Paragraph <number>"
    paragraph_prefix_pattern = re.compile(r'^Paragraph \d+:', re.IGNORECASE)

    for line in slides.strip().splitlines():
        line = line.strip()  # Remove leading and trailing whitespace
        
        if not line:  # Skip empty lines
            continue

        if line.startswith('**'):
            # Remove '**' from the title
            title = line.replace('**', '').strip()
            
            # Remove the slide prefix if it matches the pattern
            title = slide_prefix_pattern.sub('', title).strip()
            
            # Remove the paragraph prefix if it matches the pattern
            title = paragraph_prefix_pattern.sub('', title).strip()

            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = title
        elif current_slide:
            # Strip leading '-' if present, then add as bullet point
            bullet_point = line.lstrip('-').strip()
            if bullet_point:  # Ensure there's content after stripping
                text_box = current_slide.placeholders[1]
                text_frame = text_box.text_frame
                p = text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0  # Set the bullet level to 0 for top-level bullets

    return prs


from pptx.util import Inches


def main(apiKey, wordDocFilePath, topicsAndMarkersFilePath, topicWithContentSegmentFilePath,
         slideContentFilePath, outputPptPath, templatePath):

    content = py_extractContentFromWordDoc(wordDocFilePath)
    if not content:
        print("Failed to extract content from the Word document.")
        return

    # Step 1: Extract topics & markers using Claude
    extractTopicsMarkersPrompt = py_generatePrompt(EXTRACT_TOPICS_MARKERS_TEMPLATE, {"content": content})
    topicsAndMarkers = makeApiCall(apiKey, extractTopicsMarkersPrompt)
    if not topicsAndMarkers:
        print("Failed to extract topics and markers.")
        return

    py_saveFile(topicsAndMarkersFilePath, topicsAndMarkers)
    topicsMarkersJson = json.loads(py_convertTextToJson(topicsAndMarkers))

    topicWithContentSegment = ""
    slides = ""

    # Step 2: Generate AI slide content
    for i, item in enumerate(topicsMarkersJson):
        topic, marker = item['topic'], item['marker']
        try:
            startPosition = py_getMarkerPosition(marker, content, topic)

            if i < len(topicsMarkersJson) - 1:
                nextMarker = topicsMarkersJson[i + 1]['marker']
                endPosition = py_getMarkerPosition(nextMarker, content, topicsMarkersJson[i + 1]['topic'])
            else:
                endPosition = len(content)

            # ✅ Ensure we have a meaningful chunk of content
            if endPosition - startPosition < 300:
                print(f"⚠️ Short segment for topic '{topic}'. Expanding fallback window.")
                endPosition = min(len(content), startPosition + 500)

            contentSegment = py_extractContentSegment(startPosition, endPosition, content)

            if not contentSegment.strip():
                print(f"⚠️ Skipping slide generation for topic '{topic}' due to empty content segment.")
                continue

            topicWithContentSegment += f"\n\n**{topic}**\n{contentSegment}"

            generateSlideContentPrompt = py_generatePrompt(
                GENERATE_SLIDE_CONTENT_TEMPLATE,
                {"topic": topic, "contentSegment": contentSegment}
            )
            slideContent = makeApiCall(apiKey, generateSlideContentPrompt)

            if slideContent:
                slides += slideContent + "\n\n"
            else:
                print(f"⚠️ Failed to generate slide content for topic: {topic}")

        except MarkerNotFoundError as e:
            print(f"❌ Error processing topic '{topic}': {str(e)}")
            continue


    py_saveFile(topicWithContentSegmentFilePath, topicWithContentSegment)
    py_saveFile(slideContentFilePath, slides)

    # Step 3: Extract figures
    figures_info = extract_figures_from_docx(wordDocFilePath)

    # Step 4: Convert slide text to structured slide objects
    text_slides = []
    current_slide = None
    for line in slides.splitlines():
        line = line.strip()
        if not line:
            continue
        if line.startswith('**'):
            if current_slide:
                text_slides.append(current_slide)
            current_slide = {"title": line.replace('**', '').strip(), "content": ""}
        elif current_slide:
            current_slide["content"] += line + "\n"
    if current_slide:
        text_slides.append(current_slide)

    # Step 5: Map figures to slides
    figure_map = decide_slide_mapping(figures_info, text_slides)

    # Step 6: Generate initial PPT
    ppt = py_generatePPT(slides, templatePath)

    # Step 7: Place figures on corresponding slides
    from pptx.util import Inches

    for mapping in figure_map:
        figure_path = mapping.get("figure_path")
        if not figure_path or not os.path.exists(figure_path):
            print(f"⚠️ Skipping invalid image path: {figure_path}")
            continue

        if mapping["own_slide"]:
            slide = ppt.slides.add_slide(ppt.slide_layouts[1])
            slide.shapes.title.text = f"Figure {mapping['figure_id']}"
            slide.shapes.add_picture(figure_path, Inches(1), Inches(2), width=Inches(6.5), height=Inches(4))
        else:
            idx = mapping.get("slide_index")
            if idx is None or idx >= len(ppt.slides):
                print(f"⚠️ Invalid or missing slide index for Fig#{mapping['figure_id']}. Creating separate slide.")
                slide = ppt.slides.add_slide(ppt.slide_layouts[1])
                slide.shapes.title.text = f"Figure {mapping['figure_id']} (Unmatched)"
                slide.shapes.add_picture(figure_path, Inches(1), Inches(2), width=Inches(6.5), height=Inches(4))
                continue

            slide = ppt.slides[idx]

            # Try to insert beside content
            content_placeholder = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    content_placeholder = shape
                    break

            if content_placeholder:
                try:
                    content_placeholder.width = Inches(4.5)
                    content_placeholder.left = Inches(0.5)
                    slide.shapes.add_picture(
                        figure_path, Inches(5.2), Inches(2), width=Inches(4), height=Inches(3)
                    )
                except Exception as e:
                    print(f"❌ Error inserting image for Fig#{mapping['figure_id']} on slide {idx}: {e}")
            else:
                slide.shapes.add_picture(
                    figure_path, Inches(1), Inches(2), width=Inches(6.5), height=Inches(4)
                )    

    # Step 8: Save final output
    ppt.save(outputPptPath)
    print(f"\n✅ Presentation saved to: {outputPptPath}")


if __name__ == "__main__":
    load_dotenv()  # Load environment variables from .env file
    
    apiKey = os.getenv('ANTHROPIC_API_KEY')
    if not apiKey:
        print("API key not found. Please set the ANTHROPIC_API_KEY environment variable.")
        sys.exit(1)

    base_dir = os.path.dirname(os.path.abspath(__file__))
    wordDocFilePath = os.path.join(base_dir, "input/doc.docx")
    wordDocTextFilePath = os.path.join(base_dir, "intermediate/worddocInText.txt")
    topicsAndMarkersFilePath = os.path.join(base_dir, "intermediate/topics_and_markers.txt")
    topicWithContentSegmentFilePath = os.path.join(base_dir, "intermediate/topic_with_content_segment.txt")
    slideContentFilePath = os.path.join(base_dir, "intermediate/slide_content.txt")
    outputPptPath = os.path.join(base_dir, "output/output_presentation.pptx")
    templatePath = os.path.join(base_dir, "template/template.pptx")
    
    main(apiKey, wordDocFilePath, topicsAndMarkersFilePath, topicWithContentSegmentFilePath, slideContentFilePath, outputPptPath, templatePath)


