import os
import json
import re
import urllib.request
import urllib.error
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from llama_cloud_services import LlamaParse
from prompt_templates import ENRICH_PRESENTATION_PROMPT, GENERATE_SLIDE_CONTENT_TEMPLATE, EXTRACT_TOPICS_MARKERS_TEMPLATE
from fuzzywuzzy import fuzz

# Load environment variables from .env file
load_dotenv()

LLAMA_CLOUD_API_KEY = os.getenv("LLAMA_CLOUD_API_KEY")
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

if not LLAMA_CLOUD_API_KEY:
    raise ValueError("LLAMA_CLOUD_API_KEY must be set in the environment.")
if not ANTHROPIC_API_KEY:
    raise ValueError("ANTHROPIC_API_KEY must be set in the environment.")

base_dir = os.path.dirname(os.path.abspath(__file__))
intermediate_dir = os.path.join(base_dir, "intermediate")
os.makedirs(intermediate_dir, exist_ok=True)

def make_api_call_gpt(api_key, content):
    import openai
    from openai.error import RateLimitError, InvalidRequestError, AuthenticationError, APIConnectionError
    openai.api_key = api_key

    try:
        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[{"role": "user", "content": content}],
            max_tokens=4096,
            temperature=0.3
        )
        return response.choices[0].message.content.strip()

    except InvalidRequestError as e:
        print("Model not available, falling back to gpt-3.5-turbo")
        try:
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": content}],
                max_tokens=4096,
                temperature=0.3
            )
            return response.choices[0].message.content.strip()
        except Exception as inner_e:
            print(f"Fallback failed: {inner_e}")
            return None

    except RateLimitError:
        print("[Rate Limit] Too many requests or quota exhausted.")
        return None

    except AuthenticationError:
        print("[Auth Error] Check your API key and permissions.")
        return None

    except APIConnectionError:
        print("[Connection Error] Network issue when connecting to OpenAI.")
        return None

    except Exception as e:
        print(f"[Unknown OpenAI API error] {e}")
        return None
    except Exception as e:
        print(f"OpenAI API error: {e}")
        return None

def is_similar(slide_title, image_filename, threshold=70):
    clean_title = slide_title.lower().replace("_", " ")
    clean_filename = image_filename.lower().replace("_", " ")
    return fuzz.partial_ratio(clean_title, clean_filename) >= threshold

def extract_document_data(doc_path):

    parser = LlamaParse(api_key=LLAMA_CLOUD_API_KEY, language="en", verbose=True)
    result = parser.parse(doc_path)

    text_documents = result.get_text_documents(split_by_page=False)
    text = "\n".join(doc.text for doc in text_documents)

    images_dir = os.path.join(base_dir, "images")
    if os.path.exists(images_dir):
        for f in os.listdir(images_dir):
            file_path = os.path.join(images_dir, f)
            if os.path.isfile(file_path):
                os.remove(file_path)
    else:
        os.makedirs(images_dir, exist_ok=True)

    try:
        images = result.get_image_documents(
            include_screenshot_images=True,
            include_object_images=True,
            image_download_dir=images_dir
        )
    except Exception as e:
        print(f"[WARNING] Failed to fetch images from LlamaParse: {e}")
        images = []

    image_paths = []
    for idx, image_doc in enumerate(images):
        print("[DEBUG] Processing image:", idx)
        image_path = image_doc.image_path
        page_index = getattr(image_doc, 'page_index', idx)
        context_text = getattr(image_doc, 'context_text', '') or ''
        figure_label = None
        caption = ''

        # Step 2: If no LLM context, check for captions above/below
        if not context_text and page_index < len(result.pages):
            page_lines = result.pages[page_index].text.split('\n')
            for i, line in enumerate(page_lines):
                if 'fig' in line.lower():
                    context_text = line.strip()
                    break
            if not context_text:
                # Step 3: Check line just above image (assume image at bottom half of page)
                for line in reversed(page_lines):
                    if line.strip():
                        context_text = line.strip()
                        break
        if not context_text and page_index > 0:
            # Step 4: Fallback - get last line from previous page
            prev_lines = result.pages[page_index - 1].text.split('\n')
            for line in reversed(prev_lines):
                if line.strip():
                    context_text = line.strip()
                    break

        print(f"[DEBUG] context_text: {context_text}")
        
        # Try to find figure label
        figure_label = None
        match = re.search(r'(fig(?:ure)?\.?\s*\d+)', context_text, re.IGNORECASE)
        if match:
            figure_label = match.group(1).replace(" ", "_").upper()
        print(f"[DEBUG] Extracted figure_label: {figure_label}")

        # Try to extract nearby caption: 1–2 sentences max
        caption_match = re.search(r'fig(?:ure)?\.?\s*\d+\.\s*(.*)', context_text, re.IGNORECASE)
        caption = caption_match.group(1).strip() if caption_match else ''
        print(f"[DEBUG] Extracted caption: {caption}")

        parts = []
        if figure_label:
            parts.append(figure_label)
        if caption:
            parts.append('_'.join(caption.split()[:6]))  # first 6 words only

        base_name = "_".join(parts) if parts else "image"
        base_name = re.sub(r'\W+', '_', base_name)  # sanitize

        new_filename = f"img_p{page_index}_{idx}_{base_name}.png"
        new_path = os.path.join(images_dir, new_filename)

        try:
            os.rename(image_path, new_path)
            image_paths.append(new_path)
            print(f"[DEBUG] Renamed image to: {new_filename}")
        except Exception as e:
            print(f"[WARNING] Failed to rename image {image_path}: {e}")

    print(f"[DEBUG] Number of images extracted: {len(image_paths)}")

    tables = []
    for page in result.pages:
        if hasattr(page, "structuredData") and page.structuredData:
            if "tables" in page.structuredData:
                tables.extend(page.structuredData["tables"])

    print(f"[DEBUG] Number of tables extracted: {len(tables)}")

    return {"text": text, "images": image_paths, "tables": tables}


def extract_topics_from_gpt(document_text):
    prompt = EXTRACT_TOPICS_MARKERS_TEMPLATE.replace("{{content}}", document_text)
    response = make_api_call_gpt(os.getenv("OPENAI_API_KEY"), prompt)

    if response is None:
        raise ValueError("GPT topic extraction failed.")

    topics = []
    lines = response.strip().split("\n")
    for i in range(0, len(lines), 2):
        if lines[i].startswith("**") and lines[i].endswith("**"):
            topic = lines[i].strip("** ")
            if i+1 < len(lines):
                sample_text = lines[i+1].strip()
                topics.append({"topic": topic, "sample_text": sample_text})
    return topics


def split_text_into_chunks(text, max_words=1500):
    words = text.split()
    chunks = []
    for i in range(0, len(words), max_words):
        chunk = " ".join(words[i:i + max_words])
        chunks.append(chunk)
    return chunks


def enrich_with_gpt(document_data):

    print(f"[DEBUG] Table data sent to GPT: {json.dumps(document_data['tables'], indent=2)}")
    print(f"[DEBUG] Image data sent to GPT: {json.dumps(document_data['images'], indent=2)}")

    topics = extract_topics_from_gpt(document_data['text'])
    print("TOPICS:")
    for topic in topics:
        print(f"  • {topic['topic']} (sample: {topic['sample_text']})")

    segments = []
    for i in range(len(topics)):
        start = document_data['text'].find(topics[i]['sample_text'])
        end = document_data['text'].find(topics[i + 1]['sample_text']) if i + 1 < len(topics) else len(document_data['text'])
        if start != -1:
            segments.append({"topic": topics[i]['topic'], "content": document_data['text'][start:end]})

    structured_response = []

    for idx, segment in enumerate(segments):
        partial_prompt = ENRICH_PRESENTATION_PROMPT.format(
            text_content=segment['content'],
            image_paths=document_data['images'],
            table_data=json.dumps(document_data['tables'], indent=2)
        )
        print(f"[INFO] Sending topic chunk {idx+1}/{len(segments)} to GPT...")
        response = make_api_call_gpt(os.getenv("OPENAI_API_KEY"), partial_prompt)
        if response:
            try:
                structured_response.extend(json.loads(response))
                print(f"[DEBUG] GPT response parsed successfully for chunk {idx+1}.")
            except json.JSONDecodeError:
                print(f"[DEBUG] Raw GPT response (chunk {idx+1}):{response}")
                print(f"[WARNING] Topic chunk {idx+1} returned invalid JSON and was skipped.")
        else:
            print(f"[WARNING] No response for topic chunk {idx+1}, skipping.")

    if not structured_response:
        print("[ERROR] GPT API returned no usable data. Please check the prompt or API response logs above.")
        return []

    return structured_response

def add_table_to_slide(slide, table_data):
    if not table_data or not isinstance(table_data, list) or not table_data[0]:
        print("No valid table data provided.")
        return

    rows, cols = len(table_data), len(table_data[0])

    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(7.0)
    height = Inches(2.5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    total_width = Inches(7.0)
    col_width = int(total_width / cols)
    for col in range(cols):
        table.columns[col].width = col_width

    for i, row in enumerate(table_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.bold = True if i == 0 else False
            cell.text_frame.paragraphs[0].alignment = 1

    print("Table added successfully!")

def create_ppt_from_gpt(claude_slides, document_data, ppt_path="generated_deck.pptx", template_path=None):
    ppt = Presentation(template_path) if template_path and os.path.exists(template_path) else Presentation()
    if template_path and len(ppt.slides) > 0:
        rId = ppt.slides._sldIdLst[0].rId
        ppt.part.drop_rel(rId)
        del ppt.slides._sldIdLst[0]

    for slide_data in claude_slides:
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])
        if slide_data.get("title") and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]

        assigned_image = None
        slide_title = slide_data.get("title", "").lower()
        if 'used_images' not in globals():
            global used_images
            used_images = set()
        for image_path in document_data["images"]:
            if image_path in used_images:
                continue
            image_filename = os.path.basename(image_path).lower()
            if is_similar(slide_title, image_filename):
                assigned_image = image_path
                used_images.add(image_path)
                break
        if assigned_image:
            slide_data["image"] = assigned_image
        # else:
        #     print("No contextual match found, skipping image assignment for this slide.")

        slide_width = ppt.slide_width
        slide_height = ppt.slide_height
        margin = Inches(0.3)
        image_width = Inches(4.5)
        image_height = Inches(3.5)
        image_left = slide_width - image_width - margin
        image_top = (slide_height - image_height) / 2

        text_width = slide_width - image_width - Inches(1.2)
        text_box_left = Inches(1.0)
        text_box_top = Inches(2.0)
        text_box_height = slide_height - Inches(2.0)

        if len(slide.placeholders) > 1:
            text_box = slide.placeholders[1]
            text_box.left = text_box_left
            text_box.top = text_box_top
            text_box.width = text_width
            text_box.height = text_box_height
            text_frame = text_box.text_frame
            text_frame.clear()
            bullet_count = len([b for b in slide_data.get("text", "").split("\n") if b.strip()])
            dynamic_font_size = max(14, min(18, int(200 / max(bullet_count, 1))))
            for idx, line in enumerate(slide_data.get("text", "").split("\n")):
                if not line.strip(): continue
                para = text_frame.add_paragraph() if idx != 0 else text_frame.paragraphs[0]
                para.text = line
                para.bullet = True
                para.font.size = Pt(dynamic_font_size)

        if slide_data.get("image") and os.path.exists(slide_data["image"]):
            from PIL import Image
            with Image.open(slide_data["image"]) as img:
                img_width, img_height = img.size
                aspect_ratio = img_width / img_height
                max_width = Inches(4.5)
                max_height = Inches(3.5)

                if aspect_ratio > 1:
                    width = min(max_width, slide_width * 0.5)
                    height = width / aspect_ratio
                else:
                    height = min(max_height, slide_height * 0.6)
                    width = height * aspect_ratio

                image_left = slide_width - width - margin
                image_top = (slide_height - height) / 2

                slide.shapes.add_picture(slide_data["image"], image_left, image_top, width=width, height=height)

        if slide_data.get("table"):
            # print(f"[DEBUG] Slide titled '{slide_data.get('title')}' contains a table with {len(slide_data['table'])} rows")
            add_table_to_slide(slide, slide_data["table"])
        # else:
            # print(f"[DEBUG] Slide titled '{slide_data.get('title')}' has NO table")

    thank_you_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
    if thank_you_slide.shapes.title:
        thank_you_slide.shapes.title.text = "Thank You!"
    if len(thank_you_slide.placeholders) > 1:
        text_box = thank_you_slide.placeholders[1]
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.add_paragraph().text = "We appreciate your attention. Looking forward to your questions!"

    ppt.save(ppt_path)
    print(f"Presentation saved as {ppt_path}")

def main():
    doc_path = os.path.join(base_dir, "input", "doc.docx")
    output_ppt_path = os.path.join(base_dir, "output", "output_presentation.pptx")
    template_path = os.path.join(base_dir, "template", "template.pptx")

    document_data = extract_document_data(doc_path)
    with open(os.path.join(intermediate_dir, "extracted_data.json"), "w", encoding="utf-8") as json_file:
        json.dump(document_data, json_file, ensure_ascii=False, indent=4)

    gpt_response = enrich_with_gpt(document_data)
    with open(os.path.join(intermediate_dir, "gpt_structured_response.json"), "w", encoding="utf-8") as json_file:
        json.dump(gpt_response, json_file, ensure_ascii=False, indent=4)

    create_ppt_from_gpt(gpt_response, document_data, output_ppt_path, template_path)

if __name__ == "__main__":
    main()
