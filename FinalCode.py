import os
import json
import re
import uuid
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from dotenv import load_dotenv
from prompt_templates import GENERATE_SLIDE_CONTENT_TEMPLATE
from figure_extractor import extract_figures_from_docx, decide_slide_mapping
from PIL import Image
from difflib import SequenceMatcher

# === UTILS ===
def py_generatePrompt(promptTemplate, vars):
    def replace_var(match):
        var_name = match.group(1)
        return str(vars.get(var_name, f"{{{{Undefined variable: {var_name}}}}}"))
    return re.sub(r'\{\{(\w+)\}\}', replace_var, promptTemplate)

def makeApiCall(apiKey, prompt):
    import urllib.request
    import ssl
    url = "https://api.anthropic.com/v1/messages"
    headers = {
        'x-api-key': apiKey,
        'anthropic-version': '2023-06-01',
        'content-type': 'application/json'
    }
    data = json.dumps({
        "model": "claude-3-5-sonnet-20241022",
        "max_tokens": 4096,
        "messages": [{"role": "user", "content": prompt}]
    }).encode('utf-8')

    req = urllib.request.Request(url, data=data, headers=headers, method='POST')
    try:
        with urllib.request.urlopen(req) as response:
            return json.loads(response.read())['content'][0]['text']
    except Exception as e:
        print(f"API Error: {e}")
        return None

def extract_ordered_content(docx_path, image_output_folder="extracted_images"):
    from lxml import etree
    os.makedirs(image_output_folder, exist_ok=True)
    doc = Document(docx_path)
    ordered = []
    image_counter = 1
    inline_shapes_processed = set()

    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            ordered.append({"type": "paragraph", "text": paragraph.text.strip()})

        drawing_elements = paragraph._element.xpath(".//*[local-name()='drawing']")
        for drawing in drawing_elements:
            blip_elems = drawing.xpath(".//*[local-name()='blip']")
            if not blip_elems:
                continue
            embed = blip_elems[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if embed in inline_shapes_processed:
                continue
            inline_shapes_processed.add(embed)

            image_part = doc.part.related_parts[embed]
            image_bytes = image_part.blob

            filename = f"image_{image_counter}_{uuid.uuid4().hex[:8]}.png"
            path = os.path.join(image_output_folder, filename)
            with open(path, 'wb') as f:
                f.write(image_bytes)

            ordered.append({"type": "image", "path": path})
            image_counter += 1

    return ordered

def create_slides_with_inline_images(apiKey, wordDocPath, templatePath, outputPath):
    content = extract_ordered_content(wordDocPath)

    # Load the template as base and remove existing slides
    prs = Presentation(templatePath)
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    slide_layout = prs.slide_layouts[1]  # or any preferred layout

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    i = 0

    skip_phrases = [
        "I cannot create meaningful slides",
        "I would need more detailed information",
        "please provide a document portion",
        "only contains a title with no actual content",
        "Sorry, but the provided document portion",
        "I'll be happy to create relevant slides"
    ]

    last_slide_title = ""
    last_slide_bullets = []

    def is_similar(title, bullets):
        title_ratio = SequenceMatcher(None, title.lower(), last_slide_title.lower()).ratio()
        bullet_overlap = len(set(bullets).intersection(set(last_slide_bullets))) / max(len(bullets), 1)
        return title_ratio > 0.85 or bullet_overlap > 0.5

    while i < len(content):
        block = content[i]
        if block['type'] == 'paragraph':
            text = block['text']
            prompt = py_generatePrompt(GENERATE_SLIDE_CONTENT_TEMPLATE, {"topic": "", "contentSegment": text})
            slide_text = makeApiCall(apiKey, prompt)
            if not slide_text or any(skip_phrase in slide_text for skip_phrase in skip_phrases):
                i += 1
                continue

            slides_text = slide_text.strip().split("\n\n")
            for chunk in slides_text:
                lines = chunk.strip().splitlines()
                if not lines:
                    continue
                title = lines[0].strip("* ") if lines[0].startswith("**") else "Untitled"
                bullets = lines[1:] if lines[0].startswith("**") else lines

                if title == "Untitled" and len(bullets) <= 1:
                    continue

                if is_similar(title, bullets):
                    continue
                last_slide_title = title
                last_slide_bullets = bullets

                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title

                textbox = slide.placeholders[1]
                tf = textbox.text_frame
                tf.clear()
                for b in bullets:
                    if b.strip():
                        tf.add_paragraph().text = b.strip("- ")

                if i + 1 < len(content) and content[i + 1]['type'] == 'image':
                    img_path = content[i + 1]['path']
                    if os.path.exists(img_path):
                        try:
                            img = Image.open(img_path)
                            img_width, img_height = img.size

                            dpi = 96
                            img_width_inches = img_width / dpi
                            img_height_inches = img_height / dpi

                            max_img_width = 4.5
                            max_img_height = 3.5
                            scale = min(max_img_width / img_width_inches, max_img_height / img_height_inches, 1)

                            final_width = Inches(img_width_inches * scale)
                            final_height = Inches(img_height_inches * scale)

                            textbox.width = slide_width - final_width - Inches(1.5)
                            textbox.left = Inches(0.5)
                            textbox.top = Inches(1.5)
                            textbox.height = slide_height - Inches(2)

                            picture_left = slide_width - final_width - Inches(0.5)
                            picture_top = (slide_height - final_height) / 2

                            slide.shapes.add_picture(img_path, picture_left, picture_top, width=final_width, height=final_height)
                        except Exception as e:
                            print(f"Failed to add inline image: {e}")
                    i += 1
        i += 1

    prs.save(outputPath)
    print(f"Presentation saved to {outputPath}")

# === MAIN ===
def main():
    load_dotenv()
    apiKey = os.getenv("ANTHROPIC_API_KEY")
    base = os.path.dirname(os.path.abspath(__file__))
    wordDoc = os.path.join(base, "input/doc.docx")
    template = os.path.join(base, "template/template.pptx")
    output = os.path.join(base, "output/output_presentation.pptx")

    create_slides_with_inline_images(apiKey, wordDoc, template, output)

if __name__ == "__main__":
    main()
